#!/usr/bin/env python3
"""Apply Thinkport Corporate Design to an existing PowerPoint deck.

- Sets theme colors to Thinkport palette (Dark Blue, Orange, Turquoise, neutrals).
- Sets all text runs to Montserrat where possible.
- Adds background from assets/backgrounds/5.svg with a dark navy overlay (like the OpenGraph banners) and the Thinkport Venitus logo to every slide.
- Typography: Montserrat is applied to all text for consistency with the brand and OG style.
- Appends five extra example slides: Agenda, Über uns, Leistungen & Schwerpunkte,
  Referenzen, Kontakt & nächste Schritte.

Usage:
  pip install -r scripts/requirements-pptx.txt   # or: uv pip install python-pptx
  python scripts/apply-cd-to-pptx.py [--source path/to/source.pptx] [--output path/to/out.pptx]

Defaults: source = Downloads/Thinkport_Firmenvorstellung.pptx, output = examples/powerpoint/thinkport-company-intro.pptx
"""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

# Theme color names in OOXML a:clrScheme (no leading a: in xpath child names)
CLR_SCHEME_NAMES = (
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
)

# Thinkport CD: map scheme slots to hex (no #)
THINKPORT_THEME = {
    "dk1": "0B2649",       # Dark Blue – primary dark
    "lt1": "FFFFFF",       # White – text on dark
    "dk2": "333333",       # Dark Gray – body text
    "lt2": "E8E8E8",       # Light gray – secondary on dark
    "accent1": "FF5722",   # Orange – CTA, highlights
    "accent2": "00BCD4",   # Turquoise – secondary accent
    "accent3": "0B2649",   # Dark Blue
    "accent4": "00A4BA",   # Turquoise darker
    "accent5": "E64C1E",   # Orange darker
    "accent6": "174E96",   # Blue medium
    "hlink": "00BCD4",
    "folHlink": "00A4BA",
}

BRAND_FONT = "Montserrat"


def set_theme_colors(prs, theme_hex: dict[str, str]) -> None:
    """Set presentation theme colors via theme part XML (python-pptx has no API for this)."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree

    if not prs.slide_masters:
        return
    slide_master = prs.slide_masters[0]
    part = slide_master.part
    theme_part = part.part_related_by(RT.THEME)
    theme = etree.fromstring(theme_part.blob)

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    clr_scheme = theme.find(".//a:clrScheme", namespaces=ns)
    if clr_scheme is None:
        return

    for name, hex_val in theme_hex.items():
        if name not in CLR_SCHEME_NAMES:
            continue
        elem = clr_scheme.find(f".//a:{name}", namespaces=ns)
        if elem is None:
            continue
        srgb = elem.find("a:srgbClr", namespaces=ns)
        if srgb is not None:
            srgb.set("val", hex_val)
        else:
            # Replace sysClr with srgbClr so our hex is used
            sys_clr = elem.find("a:sysClr", namespaces=ns)
            if sys_clr is not None:
                elem.remove(sys_clr)
                child = etree.SubElement(elem, f"{{{ns['a']}}}srgbClr")
                child.set("val", hex_val)
                child.tail = "\n    "

    theme_part._blob = etree.tostring(theme, encoding="unicode").encode("utf-8")


def set_font_on_shape(shape, font_name: str) -> None:
    """Set font on all text runs in a shape (and in grouped shapes)."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = font_name
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name
    try:
        if shape.shapes:
            for child in shape.shapes:
                set_font_on_shape(child, font_name)
    except Exception:
        pass


def apply_font_to_slides(prs, font_name: str) -> None:
    """Set font on every text run in the presentation (slides only, not masters)."""
    for slide in prs.slides:
        for shape in slide.shapes:
            set_font_on_shape(shape, font_name)


def _ensure_background_png(root: Path, svg_path: Path, out_path: Path, width_px: int = 1920, height_px: int = 1080) -> Path:
    """Convert SVG to PNG; use cache if present and newer than SVG."""
    if out_path.exists() and out_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return out_path
    try:
        import cairosvg
        cairosvg.svg2png(
            url=str(svg_path),
            write_to=str(out_path),
            output_width=width_px,
            output_height=height_px,
        )
        return out_path
    except Exception as e:
        if out_path.exists():
            return out_path
        raise SystemExit(f"Could not convert {svg_path} to PNG (install cairosvg?). {e}")


def _background_with_navy_overlay(raw_bg_path: Path, out_path: Path, overlay_opacity: float = 0.65) -> Path:
    """Composite a dark navy overlay on the background PNG (like OpenGraph banners).
    OG uses #0B2649 at ~0.72/0.5/0 so text and logo stay readable."""
    try:
        from PIL import Image
    except ImportError:
        return raw_bg_path
    if not raw_bg_path.exists():
        return raw_bg_path
    img = Image.open(raw_bg_path).convert("RGBA")
    w, h = img.size
    overlay = Image.new("RGBA", (w, h), (11, 38, 73, int(255 * overlay_opacity)))
    out = Image.alpha_composite(img, overlay)
    out.convert("RGB").save(out_path, "PNG")
    return out_path


def _resolve_logo_path(root: Path, ppt_dir: Path) -> Path | None:
    """Return path to Venitus logo PNG (light variant for dark background)."""
    assets_dir = root / "assets"
    logo_path = assets_dir / "logos" / "venitus" / "thinkport-venitus-light.png"
    if logo_path.exists():
        return logo_path
    svg_path = assets_dir / "logos" / "venitus" / "thinkport-venitus-light.svg"
    if svg_path.exists():
        logo_png = ppt_dir / "thinkport-venitus-light.png"
        try:
            import cairosvg
            cairosvg.svg2png(url=str(svg_path), write_to=str(logo_png), output_width=400)
            return logo_png
        except Exception:
            pass
    return None


def apply_background_and_logo_to_slides(prs, root: Path) -> None:
    """Add assets/backgrounds/5.svg (as PNG) and Venitus logo to every slide.
    python-pptx MasterShapes does not support add_picture, so we add to each slide."""
    assets_dir = root / "assets"
    ppt_dir = root / "examples" / "powerpoint"
    ppt_dir.mkdir(parents=True, exist_ok=True)

    # Background: 5.svg → PNG, then add navy overlay (like OpenGraph banners) for readable text
    svg_bg = assets_dir / "backgrounds" / "5.svg"
    raw_bg = ppt_dir / "background-5.png"
    bg_png = ppt_dir / "background-5-overlay.png"
    if not svg_bg.exists():
        bg_png = None
    else:
        _ensure_background_png(root, svg_bg, raw_bg)
        _background_with_navy_overlay(raw_bg, bg_png, overlay_opacity=0.65)

    logo_path = _resolve_logo_path(root, ppt_dir)

    emu_per_inch = 914400
    logo_w_emu = int(2.2 * emu_per_inch)
    logo_h_emu = int(0.55 * emu_per_inch)
    margin_emu = int(0.35 * emu_per_inch)
    logo_left = prs.slide_width - logo_w_emu - margin_emu
    logo_top = prs.slide_height - logo_h_emu - margin_emu

    for slide in prs.slides:
        sp_tree = slide.shapes._spTree
        if bg_png and bg_png.exists():
            bg_pic = slide.shapes.add_picture(
                str(bg_png),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            sp_tree.remove(bg_pic._element)
            sp_tree.insert(0, bg_pic._element)
        if logo_path and logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path),
                logo_left,
                logo_top,
                width=logo_w_emu,
                height=logo_h_emu,
            )


# Extra example slides (German content)
EXTRA_SLIDES = [
    ("Agenda", ["• Über Thinkport", "• Leistungen & Schwerpunkte", "• Referenzen", "• Kontakt & nächste Schritte"]),
    ("Über uns", [
        "Thinkport ist spezialisiert auf Cloud, Data und KI – von der Strategie bis zur Umsetzung.",
        "",
        "• Fokus auf Microsoft Azure und moderne Datenplattformen",
        "• Starke Kombination aus Beratung und Implementierung",
        "• Langfristige Partnerschaften mit unseren Kund:innen",
    ]),
    ("Leistungen & Schwerpunkte", [
        "Cloud & Infrastruktur",
        "• Azure-Architekturen, Migration, Betrieb",
        "",
        "Daten & Analytics",
        "• Data Lakes, BI, KI/ML Use Cases",
        "",
        "Enablement & Co-Creation",
        "• Workshops, Trainings, gemeinsame Umsetzung",
    ]),
    ("Referenzen", [
        "Ausgewählte Projekte (anonymisiert):",
        "",
        "• Branchenübergreifende Cloud- und Data-Transformationen",
        "• KI-Pilotprojekte und Skalierung in die Produktion",
        "• Moderne Arbeitsweisen und Plattformen für Daten-Teams",
    ]),
    ("Kontakt & nächste Schritte", [
        "Wir freuen uns auf den Austausch.",
        "",
        "• Website: thinkport.digital",
        "• E-Mail und Ansprechpartner:innen über die Website",
        "• Nächster Schritt: kurzes Kennenlerngespräch oder Workshop",
    ]),
]


def add_example_slides(prs) -> None:
    """Append extra example slides (Agenda, Über uns, Leistungen, Referenzen, Kontakt)."""
    from pptx.util import Inches, Pt

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for title, bullets in EXTRA_SLIDES:
        slide = prs.slides.add_slide(blank_layout)
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = BRAND_FONT
        p.font.color.rgb = __rgb(0x0B, 0x26, 0x49)
        # Body
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12), Inches(5.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets):
            if i:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(12)
            p.font.name = BRAND_FONT
            p.font.color.rgb = __rgb(0x33, 0x33, 0x33)
    return


def __rgb(r: int, g: int, b: int):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    default_source = root / "examples" / "powerpoint" / "content" / "Thinkport_Firmenvorstellung.pptx"
    if not default_source.exists():
        default_source = Path.home() / "Downloads" / "Thinkport_Firmenvorstellung.pptx"
    default_out = root / "examples" / "powerpoint" / "thinkport-company-intro.pptx"

    parser = argparse.ArgumentParser(description="Apply Thinkport CD to a PowerPoint deck")
    parser.add_argument("--source", type=Path, default=default_source, help="Source PPTX path")
    parser.add_argument("--output", type=Path, default=default_out, help="Output PPTX path")
    args = parser.parse_args()

    source = args.source.resolve()
    output = args.output.resolve()
    if not source.exists():
        raise SystemExit(f"Source file not found: {source}")

    try:
        from pptx import Presentation
    except ImportError as e:
        err = f"Install python-pptx first: pip install -r scripts/requirements-pptx.txt\n{e}"
        try:
            r = Path(__file__).resolve().parent.parent
            p = r / "examples" / "powerpoint" / "apply-cd-log.txt"
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(err, encoding="utf-8")
        except Exception:
            pass
        raise SystemExit(err)

    # Copy source to output if different, then open output (so we don't mutate in place unless same path)
    if source != output:
        output.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, output)

    prs = Presentation(str(output))
    root = Path(__file__).resolve().parent.parent
    set_theme_colors(prs, THINKPORT_THEME)
    apply_font_to_slides(prs, BRAND_FONT)
    add_example_slides(prs)
    apply_background_and_logo_to_slides(prs, root)
    # Apply Montserrat again so all text (including placeholders) uses brand font (OG-style typography)
    apply_font_to_slides(prs, BRAND_FONT)
    prs.save(str(output))
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
